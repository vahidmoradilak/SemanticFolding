"""Build defense slides (PPTX, Persian RTL-ish) -> docs/thesis/fa/defense_slides.pptx"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
FIG = ROOT / "docs" / "thesis" / "fa" / "figures"
OUT = ROOT / "docs" / "thesis" / "fa" / "defense_slides.pptx"

DARK = RGBColor(0x1F, 0x3B, 0x63)
ACC = RGBColor(0xC0, 0x5A, 0x2B)

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_slide(title=None):
    s = prs.slides.add_slide(BLANK)
    if title:
        tb = s.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(12.5), Inches(0.9))
        tf = tb.text_frame; tf.word_wrap = True
        r = tf.paragraphs[0].add_run(); r.text = title
        r.font.size = Pt(30); r.font.bold = True; r.font.color.rgb = DARK
        r.font.name = "B Nazanin"
        tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
    return s


def bullets(slide, items, left=0.6, top=1.3, width=12.1, height=5.8, size=18):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for it in items:
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        lvl, txt = (it if isinstance(it, tuple) else (0, it))
        para.level = lvl
        para.alignment = PP_ALIGN.RIGHT
        bold = txt.startswith("**") and txt.endswith("**")
        clean = txt.strip("*")
        r = para.add_run(); r.text = ("• " if lvl == 0 else "– ") + clean
        r.font.size = Pt(size - lvl * 3); r.font.name = "B Nazanin"
        r.font.bold = bold or None
        r.font.color.rgb = DARK if not bold else ACC


def picture(slide, name, left=0.9, top=1.35, width=11.6):
    slide.shapes.add_picture(str(FIG / name), Inches(left), Inches(top), width=Inches(width))


def caption(slide, text, top=6.75):
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(top), Inches(12.1), Inches(0.6))
    p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(14); r.font.color.rgb = ACC; r.font.name = "B Nazanin"


# 1 ── title
s = prs.slides.add_slide(BLANK)
tb = s.shapes.add_textbox(Inches(0.7), Inches(1.9), Inches(12), Inches(2.4))
tf = tb.text_frame; tf.word_wrap = True
for i, (txt, sz) in enumerate([
    ("بررسی، ارزیابی و تحلیل کاربردهای رویکرد مبتنی بر خوشه‌بندی معنایی", 32),
    ("متون چند زبانه جهت پاسخگویی به سوالات در محیط‌های با دامنه بسته", 26),
]):
    para = tf.paragraphs[i] if i < len(tf.paragraphs) else tf.add_paragraph()
    para.alignment = PP_ALIGN.CENTER
    r = para.add_run(); r.text = txt; r.font.size = Pt(sz)
    r.font.bold = (i == 0); r.font.color.rgb = DARK; r.font.name = "B Nazanin"
tb2 = s.shapes.add_textbox(Inches(0.7), Inches(5.2), Inches(12), Inches(1.4))
tf2 = tb2.text_frame
for i, txt in enumerate(["وحید مرادی سرکشتی  |  استاد راهنما: دکتر مسعود رهگذر",
                          "دانشکده مهندسی برق و کامپیوتر، دانشگاه تهران"]):
    para = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
    para.alignment = PP_ALIGN.CENTER
    r = para.add_run(); r.text = txt; r.font.size = Pt(18); r.font.color.rgb = ACC
    r.font.name = "B Nazanin"

# 2 ── مسئله و انگیزه
s = add_slide("مسئله و انگیزه پژوهش")
bullets(s, [
    "**پرسش:** بازیابی دقیق پاسخ در حوزه‌های تخصصی (دینی، پزشکی، حقوقی) با منابع محدود",
    "مدل‌های عصبی بزرگ: دقت بالا اما نیازمند GPU، داده عظیم و برچسب‌گذاری گران",
    "الهام از مغز: HTM هاوکینز → بازنمایی پراکنده توزیع‌شده (SDR)",
    "**راه حل:** خوشه‌بندی معنایی (Semantic Folding) — نقشه معنایی دوبعدی + اثرانگشت معنایی",
    "ویژگی کلیدی برای متون دینی: استقلال زبانی اثرانگشت‌ها (عربی/انگلیسی هم‌معنا)",
])

# 3 ── roadmap
s = add_slide("نقشه‌راه هشت‌مرحله‌ای (مطابق پروپوزال)")
picture(s, "fig_roadmap.png")
caption(s, "مراحل ۱–۶ آفلاین؛ مرحله ۷ آنلاین؛ مرحله ۸ ارزیابی مستقل")

# 4 ── pipeline
s = add_slide("معماری پیاده‌سازی‌شده")
picture(s, "fig_pipeline.png")
caption(s, "نمایه‌سازی یک‌باره (S1-S5) ← پاسخگویی سریع روی SDR های تنک")

# 5 ── method core
s = add_slide("هسته روش: ساختار × اهمیت")
bullets(s, [
    "**هم‌رخدادی واژه–مفهوم = ساختار:** کدام سلول‌های نقشه روشن شوند",
    "**TF-IDF = اهمیت:** وزن = TF سند × IDF (تغذیه از یک فایل مشترک سند/پرسش)",
    "اثرانگشت سند = Σ (TF×IDF) × اثرانگشت عبارات  →  تنک‌سازی توپولوژی‌محور (10%)",
    "امتیاز: dot(q,d)/(‖q‖·√nnz) + انتشار فعال‌سازی (r=1, decay=0.5)",
    "فیوژن خطی SPLADE: s = α·SF + (1−α)·SPLADE ، α*=0.3",
], size=17)

# 6 ── datasets
s = add_slide("مجموعه‌داده‌های ارزیابی")
bullets(s, [
    "**قرآن کریم:** ۶٬۲۳۶ آیه، ۳۰ پرسش چندمرتبط — کاملاً تخصصی",
    "**دوزبانه custom_ar_en / mixed_ar_en:** ۴۸۸ قطعه موازی عربی|انگلیسی، ۴۸۸ پرسش",
    "**انگلیسی عمومی:** Belebele، NarrativeQA، PubMedQA، PopQA، MuSiQue",
    "**BEIR:** SciFact، nfcorpus، SciDocs",
    "جایگزین‌های دامنه بسته به‌جای SQuAD/NQ/XQuAD بازدامنه (توجیه در متن)",
], size=17)

# 7 ── main results
s = add_slide("نتایج اصلی: بهترین SF در برابر BM25")
picture(s, "fig_results_mrr.png", width=11.9)

# 8 ── bilingual significance
s = add_slide("دوزبانه: برتری معنادار آماری (۴۸۸ پرسش)")
bullets(s, [
    "Pure SF: MRR=0.817 > BM25=0.785   (Wilcoxon p=0.039 ، McNemar p=0.036)",
    "SF+SPLADE Linear α=0.3: MRR=0.825   (p=0.016 / 0.038)",
    "SPLADE تنها: MRR=0.481 — افت سنگین (p<1e-23)",
    "RRF: MRR=0.700 — افت معنادار (p<1e-4)",
    "**نتیجه:** توپولوژی معنایی در متون دوزبانه از تطابق واژگانی قوی‌تر است",
], size=18)

# 9 ── Quran
s = add_slide("قرآن کریم: برتری همه‌جانبه")
bullets(s, [
    "SF+SPLADE RRF در هر ۶ معیار جلوتر از BM25: MRR ×۲.۳۱ ، AP ×۳.۰۲",
    "(MRR 0.358 vs 0.155 | AP 0.218 vs 0.072)",
    "برنده در ۲۱ از ۳۰ پرسش",
    "شکست‌های باقی‌مانده: پرسش‌های موضوعی عام + نبود stemming صرفی",
], size=19)

# 10 ── ablation
s = add_slide("تحلیل پارامتری: سهم اجزا (ablation)")
picture(s, "fig_ablation_ap.png", width=11.3)
caption(s, "انتشار فعال‌سازی، تنک‌سازی متعادل و IDF حیاتی‌اند؛ گرید 64×64 عامل جهش اصلی")

# 11 ── alpha sweep
s = add_slide("فیوژن خطی: حساسیت به وزن SPLADE")
picture(s, "fig_alpha_sweep.png", width=9.6)
caption(s, "بازه بهینه α∈[0.25,0.30]؛ هم‌تراز شدن تقریباً کامل با BM25 روی Belebele")

# 12 ── neural comparison
s = add_slide("مقایسه با مدل‌های عصبی")
bullets(s, [
    "**SciFact (اجرا مستقیم):** SF+SPLADE 0.966 > BM25 0.947 > MiniLM-L6-v2 0.865",
    "SF بدون هیچ داده آموزش خارج از پیکره، از dense neural هم جلوتر زد",
    "ادبی (SemFold ارشد): BERT دقیق‌تر اما ~۲۰ برابر کندتر؛ SemFold هم‌سطح GloVe/FastText",
    "SPLADE خود یک مدل عصبی است: هم خط مبنا، هم مؤلفه فیوژن",
], size=18)

# 13 ── WordNet lesson
s = add_slide("آزمایش شبکه واژگان (WordNet): دو نسل، دو درس")
bullets(s, [
    "**v1 جانشینی سطح-متن (γ=1):** افت معنادار MRR (p=0.024)",
    "علت: بیشترِ «OOV» فقط صورت جمع بود (borders→border) که لم‌سازی داخلی هندل می‌کرد",
    "**v2 ادغام سطح-اثرانگشت (γ≤0.3) پس از لم‌سازی:** بی‌ضرر (p=0.53)، OOV حقیقی فقط ۱۴۹ پرسش",
    "درس طراحی: ادغام میرا در سطح اثرانگشت؛ سود اصلی در پیکره‌های با OOV واقعی بالا",
], size=17)

# 14 ── failures
s = add_slide("مرزها و شکست‌ها")
bullets(s, [
    "cross-lingual محض بدون پل واژگانی: MRR=0.02 → راهکار: پیکره موازی دوزبانه",
    "multi-hop (MuSiQue): نیاز به سازوکار استدلال — کار آینده GAT",
    "PopQA موجودیت‌محور: BM25 قوی‌تر می‌ماند",
    "جهش اعداد انگلیسی تا حد زیادی از بزرگ‌شدن pool کاندید (top_k→100)",
], size=19)

# 15 ── conclusions
s = add_slide("جمع‌بندی دستاوردها")
bullets(s, [
    "خط لوله بازتولیدپذیر + چارچوب بنچمارک ۱۲ مجموعه‌داده‌ای",
    "MRR=1.000 در Belebele و NarrativeQA؛ برتری معنادار دوزبانه؛ قرآن ×۲–۳",
    "جدول ablation + آزمون‌های Wilcoxon/McNemar + تحلیل شکست",
    "تعیین مرز طراحی شبکه واژگان با دو آزمایش کنترل‌شده",
], size=19)

# 16 ── future work
s = add_slide("کارهای آینده")
bullets(s, [
    "گراف توجه (GAT) برای استدلال چند-hop و گراف دانش خودکار",
    "شبکه واژگان بین‌زبانی (BabelNet) برای cross-lingual واقعی",
    "ریشه‌یابی عربی/انگلیسی؛ ارتباط معیاری؛ به‌روزرسانی پویا نقشه",
    "گسترش مقایسه تراکمی به مدل‌های بزرگ‌تر (E5/mpnet)",
], size=19)

# 17 ── thanks
s = prs.slides.add_slide(BLANK)
tb = s.shapes.add_textbox(Inches(1), Inches(3), Inches(11.3), Inches(1.5))
p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "با سپاس از توجه شما  —  پرسش‌ها؟"
r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = DARK; r.font.name = "B Nazanin"

prs.save(OUT)
print("saved ->", OUT, "| slides:", len(prs.slides.__iter__.__self__._sldIdLst))
